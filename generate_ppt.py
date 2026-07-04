import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Initialize presentation
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Design System Colors
    BG_COLOR = RGBColor(18, 18, 18)       # Slate Black (#121212)
    CARD_COLOR = RGBColor(24, 24, 24)     # Dark Gray (#181818)
    GREEN_COLOR = RGBColor(29, 185, 84)    # Spotify Green (#1DB954)
    TEXT_WHITE = RGBColor(255, 255, 255)  # White
    TEXT_GRAY = RGBColor(179, 179, 179)   # Light Gray (#B3B3B3)
    TEXT_VIOLET = RGBColor(121, 40, 202)  # Violet accent
    
    # Image Paths
    cover_img_path = r"C:\Users\sgt17\.gemini\antigravity-ide\brain\2c2d707d-e560-4883-8e9e-77dcceab626a\spotify_deck_cover_1783146449035.png"
    pm_app_img_path = r"C:\Users\sgt17\.gemini\antigravity-ide\brain\2c2d707d-e560-4883-8e9e-77dcceab626a\pm_discovery_engine_loaded_1783142747369.png"
    mvp_app_img_path = r"C:\Users\sgt17\.gemini\antigravity-ide\brain\2c2d707d-e560-4883-8e9e-77dcceab626a\discovery_queue_results_1783141230487.png"

    # Helper function to create blank slide with dark background
    def add_blank_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Draw full screen background shape for dark theme
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.color.rgb = BG_COLOR
        return slide

    # Helper to add standard slide header
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Outfit'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Add accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(1.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GREEN_COLOR
        line.line.color.rgb = GREEN_COLOR

    # Helper to add a card shape
    def add_card(slide, left, top, width, height):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = RGBColor(40, 40, 40)
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: COVER SLIDE
    # ==========================================
    slide1 = add_blank_slide()
    
    # Insert AI cover graphic (covers left half)
    if os.path.exists(cover_img_path):
        slide1.shapes.add_picture(cover_img_path, Inches(0), Inches(0), width=Inches(6.0), height=Inches(7.5))
    
    # Title box (on right half)
    title_box = slide1.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6.333), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = "Escaping the Loop:\nRe-imagining Spotify Music Discovery"
    p_title.font.name = 'Outfit'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(20)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "Building the AI-Native Discovery Co-Pilot to Reduce Repetitive Listening Behavior"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = GREEN_COLOR
    p_sub.space_after = Pt(40)
    
    # Section
    p_team = tf.add_paragraph()
    p_team.text = "Prepared by: Growth Team Product Manager\nSupporting Artifacts Hyperlinked Inside"
    p_team.font.name = 'Inter'
    p_team.font.size = Pt(14)
    p_team.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ==========================================
    slide2 = add_blank_slide()
    add_slide_header(slide2, "Aesthetic Fatigue Drives Premium Churn and Stagnates Engagement")
    
    # Left Card: The Discovery Paradox
    add_card(slide2, Inches(0.5), Inches(1.8), Inches(5.9), Inches(5.0))
    tb_left = slide2.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.3), Inches(4.4))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "The Discovery Paradox"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    bullets = [
        "A significant share of streams originates from repeat playlists, familiar artists, and previously discovered songs.",
        "Hyper-personalization is generating closed 'filter bubbles' where algorithms play safe, identical subsets of user libraries.",
        "Users actively seek discovery but retreat to repetitive habits due to high curation friction and black-box layouts.",
        "Solving this unlocks organic virality (sharing new finds) and mitigates customer retention threats."
    ]
    for b in bullets:
        p = tf_left.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    # Right Top Card: The Retention Threat
    add_card(slide2, Inches(6.9), Inches(1.8), Inches(5.9), Inches(2.3))
    tb_rt = slide2.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.3), Inches(1.9))
    tf_rt = tb_rt.text_frame
    tf_rt.word_wrap = True
    
    p = tf_rt.paragraphs[0]
    p.text = "The Churn Threat"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    p = tf_rt.add_paragraph()
    p.text = "• Daily Premium listeners experience 'aesthetic boredom' from repetitive loops, motivating them to switch to competitors (Apple Music, YouTube Music) or search manually on Bandcamp/SoundCloud."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY

    # Right Bottom Card: The Core Opportunity
    add_card(slide2, Inches(6.9), Inches(4.5), Inches(5.9), Inches(2.3))
    tb_rb = slide2.shapes.add_textbox(Inches(7.2), Inches(4.7), Inches(5.3), Inches(1.9))
    tf_rb = tb_rb.text_frame
    tf_rb.word_wrap = True
    
    p = tf_rb.paragraphs[0]
    p.text = "The AI-Native Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    p = tf_rb.add_paragraph()
    p.text = "• Move away from passive recommendation structures. Give users interactive inputs (Discovery Intensity dial and Niche Boost filters) back-ended by LLM-semantic matching to rebuild platform trust."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 3: REVIEW INSIGHTS
    # ==========================================
    slide3 = add_blank_slide()
    add_slide_header(slide3, "User Reviews Reveal Polarized Sentiment Driven by Shuffle Loops and Ad Friction")
    
    # Metrics columns
    add_card(slide3, Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.8))
    tb_met = slide3.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(3.4), Inches(4.2))
    tf_met = tb_met.text_frame
    tf_met.word_wrap = True
    
    p = tf_met.paragraphs[0]
    p.text = "Feedback Overview"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(20)
    
    p = tf_met.add_paragraph()
    p.text = "55%"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_met.add_paragraph()
    p.text = "Positive Sentiment (Vast library catalog, concept satisfaction)"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    p.space_after = Pt(20)
    
    p = tf_met.add_paragraph()
    p.text = "35%"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_met.add_paragraph()
    p.text = "Negative Sentiment (Shuffle repetition, smart shuffle bias, ads)"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    
    # Details Card
    add_card(slide3, Inches(4.8), Inches(1.8), Inches(8.0), Inches(4.8))
    tb_det = slide3.shapes.add_textbox(Inches(5.1), Inches(2.1), Inches(7.4), Inches(4.2))
    tf_det = tb_det.text_frame
    tf_det.word_wrap = True
    
    p = tf_det.paragraphs[0]
    p.text = "Dominant Review Themes & Core Complaints"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    
    complaints = [
        "The Shuffle Loop: Users report that even with 1,000+ song playlists, the shuffle feature cycles through the 'same fifty songs' repeatedly, ignoring saved gems.",
        "Smart Shuffle Intrusion: Free and paying users criticize Smart Shuffle for injecting unwanted, repetitive pop songs that clash with the playlist's specific vibe.",
        "UI Overload: Intrusion of non-music content (podcasts, audiobooks) on the home screen increases browsing friction and drives users to loop cached favorites.",
        "Review Analysis Workflow: Supported by our dynamic analysis workspace."
    ]
    for c in complaints:
        p = tf_det.add_paragraph()
        p.text = "•  " + c
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(8)
        
    # Hyperlink to report
    p_link = tf_det.add_paragraph()
    p_link.space_before = Pt(15)
    p_link.text = "Read Full Feedback Database and PM Report: "
    p_link.font.size = Pt(14)
    p_link.font.color.rgb = TEXT_WHITE
    run = p_link.add_run()
    run.text = "analysis_report.md"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN_COLOR
    run.hyperlink.address = "file:///c:/Users/sgt17/Downloads/NL_Spotify/data/analysis_report.md"

    # ==========================================
    # SLIDE 4: PERSONAS
    # ==========================================
    slide4 = add_blank_slide()
    add_slide_header(slide4, "Curators and Specialists Face Different Discovery Barriers")
    
    # 3 Cards side-by-side
    card_w = Inches(3.8)
    card_h = Inches(4.6)
    card_y = Inches(1.8)
    
    # Persona 1: Sarah
    add_card(slide4, Inches(0.5), card_y, card_w, card_h)
    tb = slide4.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(3.4), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sarah, 28\nThe Playlist Curator"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    bullets = [
        "Manages massive custom playlists (1,500+ tracks).",
        "Frustration: Shuffle loops on a tiny fraction of songs. Smart Shuffle injects commercial pop that ruins the vibe.",
        "Workaround: Uses Shazam in public or scans external blogs, then manually inputs songs to Spotify."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    # Persona 2: Marcus
    add_card(slide4, Inches(4.7), card_y, card_w, card_h)
    tb2 = slide4.shapes.add_textbox(Inches(4.9), Inches(2.0), Inches(3.4), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Marcus, 34\nThe Genre Specialist"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    
    bullets = [
        "Listens to niche genres (ambient electronica, indie jazz).",
        "Frustration: Spotify recommends generic, major-label crossover pop. AI DJ transitions are jarring and disrupt flow.",
        "Workaround: Discovers new music on Bandcamp/SoundCloud, then streams on Spotify only as player."
    ]
    for b in bullets:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    # Persona 3: Priya
    add_card(slide4, Inches(8.9), card_y, card_w, card_h)
    tb3 = slide4.shapes.add_textbox(Inches(9.1), Inches(2.0), Inches(3.4), Inches(4.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Priya, 22\nThe Casual Discoverer"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    
    bullets = [
        "Relies on background music (Daily Mixes) for focus.",
        "Frustration: Daily Mixes become static loops repeating liked songs instead of suggesting new releases.",
        "Workaround: Overcome by decision fatigue, she accepts repetitive streams because active discovery is painful."
    ]
    for b in bullets:
        p = tf3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    # Link at bottom
    tb_link = slide4.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf_link = tb_link.text_frame
    p_link = tf_link.paragraphs[0]
    p_link.text = "Access Primary Research and Interview Transcripts: "
    p_link.font.size = Pt(14)
    p_link.font.color.rgb = TEXT_WHITE
    run = p_link.add_run()
    run.text = "user_research_validation.md"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN_COLOR
    run.hyperlink.address = "file:///C:/Users/sgt17/.gemini/antigravity-ide/brain/2c2d707d-e560-4883-8e9e-77dcceab626a/user_research_validation.md"

    # ==========================================
    # SLIDE 5: ROOT CAUSES
    # ==========================================
    slide5 = add_blank_slide()
    add_slide_header(slide5, "Algorithmic Feedback Loops and UX Clutter Root Cause Stagnation")
    
    add_card(slide5, Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.8))
    tb_rc1 = slide5.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(3.4), Inches(4.2))
    tf_rc1 = tb_rc1.text_frame
    tf_rc1.word_wrap = True
    p = tf_rc1.paragraphs[0]
    p.text = "1. Technical Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    bullets = [
        "Collaborative Filtering Lock-in: Matrix factorization treats historical listens as permanent preference, trapping user vectors in tight local sub-clusters.",
        "Shuffle Optimization: Uses local caching and affinity bias instead of true randomness to reduce streaming server requests."
    ]
    for b in bullets:
        p = tf_rc1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    add_card(slide5, Inches(4.7), Inches(1.8), Inches(3.8), Inches(4.8))
    tb_rc2 = slide5.shapes.add_textbox(Inches(4.9), Inches(2.1), Inches(3.4), Inches(4.2))
    tf_rc2 = tb_rc2.text_frame
    tf_rc2.word_wrap = True
    p = tf_rc2.paragraphs[0]
    p.text = "2. UX / Design Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    bullets = [
        "The Black-Box Interface: Recommendations offer zero sliders to modify the novelty vs. familiarity ratio.",
        "Home Screen Clutter: High-visibility rows are loaded with podcasts and audiobooks, increasing search cost and driving users to loop downloaded tracks."
    ]
    for b in bullets:
        p = tf_rc2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    add_card(slide5, Inches(8.9), Inches(1.8), Inches(3.8), Inches(4.8))
    tb_rc3 = slide5.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.4), Inches(4.2))
    tf_rc3 = tb_rc3.text_frame
    tf_rc3.word_wrap = True
    p = tf_rc3.paragraphs[0]
    p.text = "3. Commercial Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    bullets = [
        "Label Payola placements: Smart Shuffle algorithms boost sponsored releases via marketing distribution deals.",
        "Royalty Cost Structure: Heavy streaming of major-label pop catalog incurs high licensing margins, disincentivizing organic long-tail discovery."
    ]
    for b in bullets:
        p = tf_rc3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 6: BUSINESS CASE
    # ==========================================
    slide6 = add_blank_slide()
    add_slide_header(slide6, "Retaining Bored Subscribers and Shifting Stream-Share Optimizes Gross Margins")
    
    # KPI 1: Churn Reduction
    add_card(slide6, Inches(0.5), Inches(1.8), Inches(5.9), Inches(2.3))
    tb_k1 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.3), Inches(1.9))
    tf_k1 = tb_k1.text_frame
    tf_k1.word_wrap = True
    p = tf_k1.paragraphs[0]
    p.text = "Protecting Premium Churn"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(5)
    p = tf_k1.add_paragraph()
    p.text = "Aesthetic fatigue is a top reason premium users churn. Reducing monthly subscriber churn by 0.1% across 240M premium users yields a run-rate increase of ~$34.5 Million in Annual Recurring Revenue (ARR)."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE

    # KPI 2: Margin Optimization
    add_card(slide6, Inches(6.9), Inches(1.8), Inches(5.9), Inches(2.3))
    tb_k2 = slide6.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.3), Inches(1.9))
    tf_k2 = tb_k2.text_frame
    tf_k2.word_wrap = True
    p = tf_k2.paragraphs[0]
    p.text = "Content Royalty Optimization"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(5)
    p = tf_k2.add_paragraph()
    p.text = "Streaming major-label pop hits carries expensive minimum payout thresholds. Distributing 400 bps of stream share to independent, long-tail artists via niche discovery features improves Spotify's gross margins on content licensing."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY

    # Narrative Box
    add_card(slide6, Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.1))
    tb_nb = slide6.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.733), Inches(1.7))
    tf_nb = tb_nb.text_frame
    tf_nb.word_wrap = True
    p = tf_nb.paragraphs[0]
    p.text = "The Strategic Loop: Trust, Engagement, and Virality"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    bullets = [
        "Restoring Recommendation Trust: Showing users why songs are recommended stops them from turning to Bandcamp or SoundCloud.",
        "Organic Growth Loop: Listeners who discover obscure local talent share tracks on social platforms, creating a low-cost organic acquisition engine."
    ]
    for b in bullets:
        p = tf_nb.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(5)

    # ==========================================
    # SLIDE 7: AI-NATIVE ARCHITECTURE
    # ==========================================
    slide7 = add_blank_slide()
    add_slide_header(slide7, "Generative AI Unlocks Semantic Matching and Multi-Constraint Controls")
    
    # 2 Columns comparing systems
    # Column 1: Traditional
    add_card(slide7, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_t = slide7.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.3), Inches(4.2))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p = tf_t.paragraphs[0]
    p.text = "Traditional Recommender Systems"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_GRAY
    p.space_after = Pt(15)
    
    items = [
        "Vector Correlation Bubble: CF matches user profile lookalikes. Recommends songs similar listeners heard, leading to high repetition.",
        "Keyword Restrictions: Search requires hard metadata matches. Fails on abstract searches (e.g. 'moody synth for rainy night in Paris').",
        "Black-Box Placement: Recommendations are static queues. No user knobs to scale novelty or filter label commercial content."
    ]
    for item in items:
        p = tf_t.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    # Column 2: AI-Native
    add_card(slide7, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_a = slide7.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.2))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    p = tf_a.paragraphs[0]
    p.text = "AI-Native Co-Pilot"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    items = [
        "Semantic Search Mapping: Multi-dimensional embedding maps abstract vibe descriptions to sonic profiles containing detailed texture descriptions.",
        "Multi-Constraint Reasoning: Interprets natural language vibe + novelty dial + Niche Boost settings simultaneously to build custom play queues.",
        "Transparent Rationales: Explains exactly why each track is recommended, restoring user trust and encouraging exploration."
    ]
    for item in items:
        p = tf_a.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 8: THE 1-SLIDER WORKFLOW
    # ==========================================
    slide8 = add_blank_slide()
    add_slide_header(slide8, "The AI-Powered Review Analysis Engine Ingests, Normalizes, and Summarizes Feedback")
    
    # Left Box: Workflow Diagram
    add_card(slide8, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_wf = slide8.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_wf = tb_wf.text_frame
    tf_wf.word_wrap = True
    p = tf_wf.paragraphs[0]
    p.text = "Review Analysis Workflow (PM tool)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    steps = [
        "1. Scraping Engines: Automated scripts scrape App Store and Google Play reviews, filtering for feedback within the last 12 months.",
        "2. Normalizer Module: Deduplicates reviews, aligns schemas, and removes entries under 10 characters to filter out noise.",
        "3. LLM Pipeline: Feeds normalized reviews to Gemini with structured PM prompts to cluster themes, personas, and JTBDs.",
        "4. Ask the Reviews Chat: Integrates a Q&A chatbot allowing product managers to query the feedback database in plain language."
    ]
    for s in steps:
        p = tf_wf.add_paragraph()
        p.text = s
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    p_link = tf_wf.add_paragraph()
    p_link.space_before = Pt(15)
    p_link.text = "Test out the Workflow: "
    p_link.font.size = Pt(14)
    p_link.font.color.rgb = TEXT_WHITE
    run = p_link.add_run()
    run.text = "PM Review Dashboard (Port 8501)"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN_COLOR
    run.hyperlink.address = "http://localhost:8501"

    # Right Box: Embed screenshot of PM dashboard
    add_card(slide8, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8))
    if os.path.exists(pm_app_img_path):
        slide8.shapes.add_picture(pm_app_img_path, Inches(7.1), Inches(2.0), width=Inches(5.5), height=Inches(4.4))

    # ==========================================
    # SLIDE 9: CO-PILOT MVP WALKTHROUGH
    # ==========================================
    slide9 = add_blank_slide()
    add_slide_header(slide9, "Our Live Standalone Co-Pilot MVP Empowers Active Novelty Tuning")
    
    # Left Box: MVP Controls Walkthrough
    add_card(slide9, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_mvp = slide9.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_mvp = tb_mvp.text_frame
    tf_mvp.word_wrap = True
    p = tf_mvp.paragraphs[0]
    p.text = "Co-Pilot Interactive Tuning Features"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    features = [
        "Vibe Description Input: Users describe a mental state or scene. LLM handles semantic cross-references against track details.",
        "Discovery Intensity Slider: Select 0% for familiar hits, 50% for balanced mixes, or 100% to explore highly obscure independent music.",
        "Niche Boost Switch: Active filter that prioritizes independent releases (< 65 popularity index) over commercial pop.",
        "Transparent Explanations: Each track displays an explanation detailing how its tempo, rhythm, or instruments fit the vibe."
    ]
    for f in features:
        p = tf_mvp.add_paragraph()
        p.text = "• " + f
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    p_link = tf_mvp.add_paragraph()
    p_link.space_before = Pt(10)
    p_link.text = "Play with the Live Prototype: "
    p_link.font.size = Pt(14)
    p_link.font.color.rgb = TEXT_WHITE
    run = p_link.add_run()
    run.text = "Discovery Co-Pilot (Port 8502)"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN_COLOR
    run.hyperlink.address = "http://localhost:8502"

    # Right Box: Embed screenshot of Discovery Co-Pilot MVP
    add_card(slide9, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8))
    if os.path.exists(mvp_app_img_path):
        slide9.shapes.add_picture(mvp_app_img_path, Inches(7.1), Inches(2.0), width=Inches(5.5), height=Inches(4.4))

    # ==========================================
    # SLIDE 10: ROLL-OUT & SUCCESS METRICS
    # ==========================================
    slide10 = add_blank_slide()
    add_slide_header(slide10, "We Will Measure the Stream Diversity Index to Validate Retention Impact")
    
    # Left Card: Target Metrics
    add_card(slide10, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_m1 = slide10.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.3), Inches(4.2))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True
    p = tf_m1.paragraphs[0]
    p.text = "Key Metrics for Success"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    p.space_after = Pt(15)
    
    metrics = [
        "Stream Diversity Index (North Star): Gini coefficient of artists played per user over a 30-day period. Target: Reduce Gini concentration by 15% (more equal share).",
        "Recommendation Save Rate: Percentage of recommended tracks saved to 'Liked Songs' or custom playlists. Target: Increase from ~3% to 8%.",
        "Discovery Share of Streams: Percentage of total streams coming from co-pilot queue vs. user library. Target: 35%.",
        "Churn Mitigation: 5% churn reduction among target cohort."
    ]
    for m in metrics:
        p = tf_m1.add_paragraph()
        p.text = "• " + m
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
    # Right Card: Release Roadmap
    add_card(slide10, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8))
    tb_m2 = slide10.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.2))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True
    p = tf_m2.paragraphs[0]
    p.text = "Deployment & Roll-out Strategy"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(15)
    
    phases = [
        "Phase 1: A/B Testing (Weeks 1-4)\nDeploy Co-Pilot to a 5% cohort of premium subscribers exhibiting high boredom signals (declining listening hours, repeated skips).",
        "Phase 2: Regional Alpha (Weeks 5-8)\nExpand release to a specific region (e.g. US West or India) to test server/LLM latency, embedding API costs, and load capacity.",
        "Phase 3: Global Beta Integration (Week 9+)\nMerge controls directly under the standard 'Search' or 'AI DJ' feed, replacing Smart Shuffle entirely."
    ]
    for p_text in phases:
        p = tf_m2.add_paragraph()
        p.text = "• " + p_text
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(8)

    # 3. Save Presentation
    output_filename = "NL Spotify.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}' with slide count {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
